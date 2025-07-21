"""
Multi-Modal Processing Services for Second Brain v2.6.0
Comprehensive services for processing text, audio, video, image, and document content
"""

import asyncio
import hashlib
import io
import json
import logging
import mimetypes
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
from uuid import UUID, uuid4

import aiofiles
import numpy as np
from PIL import Image, ExifTags
from pydantic import BaseModel

# Audio processing
try:
    import librosa
    import soundfile as sf
    import whisper
except ImportError:
    librosa = sf = whisper = None

# Video processing  
try:
    import cv2
    import moviepy.editor as mp
except ImportError:
    cv2 = mp = None

# Document processing
try:
    import docx
    import openpyxl
    import pypdf2
    from pptx import Presentation
except ImportError:
    docx = openpyxl = pypdf2 = Presentation = None

# OCR
try:
    import pytesseract
    import easyocr
except ImportError:
    pytesseract = easyocr = None

# AI models
try:
    import openai
    import clip
    import torch
except ImportError:
    openai = clip = torch = None

from .models import (
    ContentType,
    ProcessingStatus,
    AudioMetadata,
    VideoMetadata,
    ImageMetadata,
    DocumentMetadata
)

logger = logging.getLogger(__name__)


class ProcessingError(Exception):
    """Custom exception for processing errors."""
    pass


class BaseProcessor:
    """Base class for all content processors."""
    
    def __init__(self):
        self.supported_formats = []
        self.max_file_size = 100 * 1024 * 1024  # 100MB default
    
    async def can_process(self, file_path: str, mime_type: str) -> bool:
        """Check if this processor can handle the given file."""
        return mime_type in self.supported_formats
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process the file and return extracted data."""
        raise NotImplementedError
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate SHA-256 hash of file."""
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    
    def _get_file_size(self, file_path: str) -> int:
        """Get file size in bytes."""
        return os.path.getsize(file_path)


class TextProcessor(BaseProcessor):
    """Processor for plain text files."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [
            'text/plain',
            'text/csv',
            'text/html',
            'text/xml',
            'application/json',
            'application/xml'
        ]
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process text file."""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            
            return {
                'content': content,
                'extracted_text': content,
                'word_count': len(content.split()),
                'character_count': len(content),
                'file_hash': self._get_file_hash(file_path),
                'file_size': self._get_file_size(file_path)
            }
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            raise ProcessingError(f"Text processing failed: {e}")


class ImageProcessor(BaseProcessor):
    """Processor for image files."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [
            'image/jpeg',
            'image/png', 
            'image/gif',
            'image/bmp',
            'image/tiff',
            'image/webp'
        ]
        # Load CLIP model if available
        self.clip_model = None
        self.clip_preprocess = None
        if clip and torch:
            try:
                self.clip_model, self.clip_preprocess = clip.load("ViT-B/32")
            except Exception as e:
                logger.warning(f"Failed to load CLIP model: {e}")
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process image file."""
        try:
            with Image.open(file_path) as img:
                # Basic image info
                metadata = {
                    'width': img.width,
                    'height': img.height,
                    'format': img.format,
                    'color_mode': img.mode,
                    'file_hash': self._get_file_hash(file_path),
                    'file_size': self._get_file_size(file_path)
                }
                
                # EXIF data
                if hasattr(img, '_getexif') and img._getexif():
                    exif = {
                        ExifTags.TAGS.get(k, k): v 
                        for k, v in img._getexif().items()
                        if k in ExifTags.TAGS
                    }
                    metadata['exif_data'] = exif
                
                # Dominant colors
                img_rgb = img.convert('RGB')
                colors = img_rgb.getcolors(maxcolors=256*256*256)
                if colors:
                    dominant_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:5]
                    metadata['dominant_colors'] = [
                        {'rgb': color[1], 'count': color[0]} 
                        for color in dominant_colors
                    ]
                
                # OCR if requested
                extracted_text = ""
                if kwargs.get('extract_text', True) and pytesseract:
                    try:
                        extracted_text = pytesseract.image_to_string(img)
                        metadata['ocr_confidence'] = pytesseract.image_to_data(
                            img, output_type=pytesseract.Output.DICT
                        )
                    except Exception as e:
                        logger.warning(f"OCR failed for {file_path}: {e}")
                
                # Scene description using CLIP if available
                scene_description = ""
                if self.clip_model and kwargs.get('analyze_content', True):
                    try:
                        scene_description = await self._generate_scene_description(img)
                    except Exception as e:
                        logger.warning(f"Scene description failed: {e}")
                
                return {
                    'content_type': ContentType.IMAGE,
                    'extracted_text': extracted_text,
                    'scene_description': scene_description,
                    'metadata': metadata
                }
                
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            raise ProcessingError(f"Image processing failed: {e}")
    
    async def _generate_scene_description(self, img: Image.Image) -> str:
        """Generate scene description using CLIP."""
        if not self.clip_model:
            return ""
        
        try:
            # Preprocess image
            image_tensor = self.clip_preprocess(img).unsqueeze(0)
            
            # Common scene descriptions to test against
            text_options = [
                "a photo of a person",
                "a photo of an animal",
                "a photo of a building",
                "a photo of nature",
                "a photo of food",
                "a photo of a vehicle",
                "a photo of text or document",
                "a photo of artwork",
                "a photo of technology",
                "a screenshot of software"
            ]
            
            text_tokens = clip.tokenize(text_options)
            
            with torch.no_grad():
                image_features = self.clip_model.encode_image(image_tensor)
                text_features = self.clip_model.encode_text(text_tokens)
                
                # Calculate similarities
                similarities = (100.0 * image_features @ text_features.T).softmax(dim=-1)
                values, indices = similarities[0].topk(3)
                
                descriptions = []
                for value, index in zip(values, indices):
                    descriptions.append(f"{text_options[index]} ({value.item():.1%})")
                
                return "; ".join(descriptions)
                
        except Exception as e:
            logger.error(f"CLIP scene description failed: {e}")
            return ""


class AudioProcessor(BaseProcessor):
    """Processor for audio files."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [
            'audio/mpeg',
            'audio/wav',
            'audio/flac',
            'audio/ogg',
            'audio/m4a',
            'audio/mp3'
        ]
        # Load Whisper model if available
        self.whisper_model = None
        if whisper:
            try:
                self.whisper_model = whisper.load_model("base")
            except Exception as e:
                logger.warning(f"Failed to load Whisper model: {e}")
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process audio file."""
        try:
            metadata = {
                'file_hash': self._get_file_hash(file_path),
                'file_size': self._get_file_size(file_path)
            }
            
            # Load audio with librosa
            if librosa:
                y, sr = librosa.load(file_path)
                duration = librosa.get_duration(y=y, sr=sr)
                
                metadata.update({
                    'duration': duration,
                    'sample_rate': sr,
                    'channels': 1 if y.ndim == 1 else y.shape[0]
                })
                
                # Audio features
                if kwargs.get('analyze_content', True):
                    try:
                        tempo, _ = librosa.beat.beat_track(y=y, sr=sr)
                        spectral_centroid = librosa.feature.spectral_centroid(y=y, sr=sr)
                        zero_crossing_rate = librosa.feature.zero_crossing_rate(y)
                        
                        metadata['audio_features'] = {
                            'tempo': float(tempo),
                            'spectral_centroid_mean': float(np.mean(spectral_centroid)),
                            'zero_crossing_rate_mean': float(np.mean(zero_crossing_rate))
                        }
                    except Exception as e:
                        logger.warning(f"Audio feature extraction failed: {e}")
            
            # Transcription
            transcript = ""
            transcript_confidence = 0.0
            if self.whisper_model and kwargs.get('extract_text', True):
                try:
                    result = self.whisper_model.transcribe(file_path)
                    transcript = result.get('text', '')
                    # Whisper doesn't provide word-level confidence, so we estimate
                    transcript_confidence = 0.8 if transcript else 0.0
                    
                    # Language detection
                    if 'language' in result:
                        metadata['language'] = result['language']
                        
                except Exception as e:
                    logger.warning(f"Audio transcription failed: {e}")
            
            return {
                'content_type': ContentType.AUDIO,
                'extracted_text': transcript,
                'metadata': {
                    **metadata,
                    'transcript': transcript,
                    'transcript_confidence': transcript_confidence
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {e}")
            raise ProcessingError(f"Audio processing failed: {e}")


class VideoProcessor(BaseProcessor):
    """Processor for video files."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [
            'video/mp4',
            'video/avi',
            'video/mov',
            'video/wmv',
            'video/flv',
            'video/webm'
        ]
        self.max_file_size = 500 * 1024 * 1024  # 500MB for videos
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process video file."""
        try:
            metadata = {
                'file_hash': self._get_file_hash(file_path),
                'file_size': self._get_file_size(file_path)
            }
            
            # Process with moviepy
            if mp:
                try:
                    with mp.VideoFileClip(file_path) as video:
                        metadata.update({
                            'duration': video.duration,
                            'width': video.w,
                            'height': video.h,
                            'fps': video.fps,
                            'has_audio': video.audio is not None
                        })
                        
                        if video.audio:
                            # Extract audio for transcription
                            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_audio:
                                video.audio.write_audiofile(tmp_audio.name, verbose=False, logger=None)
                                
                                # Process audio for transcription
                                audio_processor = AudioProcessor()
                                if await audio_processor.can_process(tmp_audio.name, 'audio/wav'):
                                    audio_result = await audio_processor.process(tmp_audio.name, **kwargs)
                                    metadata['transcript'] = audio_result['extracted_text']
                                
                                # Clean up temp file
                                os.unlink(tmp_audio.name)
                        
                        # Extract keyframes
                        if kwargs.get('analyze_content', True):
                            try:
                                keyframes = []
                                duration = video.duration
                                # Extract frames at regular intervals
                                for t in np.linspace(0, duration, min(10, int(duration))):
                                    frame = video.get_frame(t)
                                    # Convert frame to PIL Image for further processing
                                    frame_img = Image.fromarray(frame)
                                    
                                    keyframes.append({
                                        'timestamp': t,
                                        'width': frame_img.width,
                                        'height': frame_img.height
                                    })
                                
                                metadata['keyframes'] = keyframes
                                metadata['frame_count'] = int(duration * video.fps)
                                
                            except Exception as e:
                                logger.warning(f"Keyframe extraction failed: {e}")
                                
                except Exception as e:
                    logger.warning(f"MoviePy processing failed: {e}")
                    # Fallback to OpenCV if available
                    if cv2:
                        await self._process_with_opencv(file_path, metadata)
            
            return {
                'content_type': ContentType.VIDEO,
                'extracted_text': metadata.get('transcript', ''),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing video {file_path}: {e}")
            raise ProcessingError(f"Video processing failed: {e}")
    
    async def _process_with_opencv(self, file_path: str, metadata: Dict[str, Any]):
        """Fallback video processing with OpenCV."""
        try:
            cap = cv2.VideoCapture(file_path)
            
            # Get video properties
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = cap.get(cv2.CAP_PROP_FRAME_COUNT)
            width = cap.get(cv2.CAP_PROP_FRAME_WIDTH)
            height = cap.get(cv2.CAP_PROP_FRAME_HEIGHT)
            
            metadata.update({
                'duration': frame_count / fps if fps > 0 else 0,
                'width': int(width),
                'height': int(height),
                'fps': fps,
                'frame_count': int(frame_count)
            })
            
            cap.release()
            
        except Exception as e:
            logger.warning(f"OpenCV processing failed: {e}")


class DocumentProcessor(BaseProcessor):
    """Processor for document files."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [
            'application/pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/msword',
            'application/vnd.ms-powerpoint',
            'application/vnd.ms-excel'
        ]
    
    async def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process document file."""
        try:
            metadata = {
                'file_hash': self._get_file_hash(file_path),
                'file_size': self._get_file_size(file_path)
            }
            
            mime_type, _ = mimetypes.guess_type(file_path)
            extracted_text = ""
            
            # PDF processing
            if mime_type == 'application/pdf' and pypdf2:
                extracted_text, pdf_metadata = await self._process_pdf(file_path)
                metadata.update(pdf_metadata)
            
            # Word document processing  
            elif 'wordprocessingml' in (mime_type or '') and docx:
                extracted_text, doc_metadata = await self._process_docx(file_path)
                metadata.update(doc_metadata)
            
            # PowerPoint processing
            elif 'presentationml' in (mime_type or '') and Presentation:
                extracted_text, ppt_metadata = await self._process_pptx(file_path)
                metadata.update(ppt_metadata)
            
            # Excel processing
            elif 'spreadsheetml' in (mime_type or '') and openpyxl:
                extracted_text, xl_metadata = await self._process_xlsx(file_path)
                metadata.update(xl_metadata)
            
            return {
                'content_type': ContentType.DOCUMENT,
                'extracted_text': extracted_text,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            raise ProcessingError(f"Document processing failed: {e}")
    
    async def _process_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PDF file."""
        text = ""
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                reader = pypdf2.PdfFileReader(file)
                
                metadata.update({
                    'page_count': reader.numPages,
                    'format': 'PDF'
                })
                
                # Extract document info
                if reader.documentInfo:
                    info = reader.documentInfo
                    metadata.update({
                        'title': info.get('/Title', ''),
                        'author': info.get('/Author', ''),
                        'subject': info.get('/Subject', ''),
                        'creator': info.get('/Creator', ''),
                        'creation_date': info.get('/CreationDate', ''),
                        'modification_date': info.get('/ModDate', '')
                    })
                
                # Extract text from all pages
                for page_num in range(reader.numPages):
                    page = reader.getPage(page_num)
                    page_text = page.extractText()
                    text += page_text + "\n\n"
                
                metadata['word_count'] = len(text.split())
                metadata['character_count'] = len(text)
                
        except Exception as e:
            logger.warning(f"PDF processing error: {e}")
        
        return text, metadata
    
    async def _process_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process Word document."""
        text = ""
        metadata = {}
        
        try:
            doc = docx.Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            text = "\n\n".join(paragraphs)
            
            # Document properties
            props = doc.core_properties
            metadata.update({
                'format': 'DOCX',
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'word_count': len(text.split()),
                'character_count': len(text),
                'creation_date': props.created,
                'modification_date': props.modified
            })
            
        except Exception as e:
            logger.warning(f"DOCX processing error: {e}")
        
        return text, metadata
    
    async def _process_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PowerPoint presentation."""
        text = ""
        metadata = {}
        
        try:
            prs = Presentation(file_path)
            
            slides_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text.append("\n".join(slide_text))
            
            text = "\n\n--- SLIDE BREAK ---\n\n".join(slides_text)
            
            metadata.update({
                'format': 'PPTX',
                'page_count': len(prs.slides),
                'word_count': len(text.split()),
                'character_count': len(text)
            })
            
        except Exception as e:
            logger.warning(f"PPTX processing error: {e}")
        
        return text, metadata
    
    async def _process_xlsx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process Excel spreadsheet."""
        text = ""
        metadata = {}
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = []
                
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):  # Skip empty rows
                        sheet_data.append("\t".join(row_text))
                
                if sheet_data:
                    sheets_text.append(f"=== {sheet_name} ===\n" + "\n".join(sheet_data))
            
            text = "\n\n".join(sheets_text)
            
            metadata.update({
                'format': 'XLSX',
                'sheet_count': len(wb.sheetnames),
                'word_count': len(text.split()),
                'character_count': len(text)
            })
            
        except Exception as e:
            logger.warning(f"XLSX processing error: {e}")
        
        return text, metadata


class MultiModalProcessingService:
    """Main service for coordinating multimodal processing."""
    
    def __init__(self):
        self.processors = {
            ContentType.TEXT: TextProcessor(),
            ContentType.IMAGE: ImageProcessor(),
            ContentType.AUDIO: AudioProcessor(),
            ContentType.VIDEO: VideoProcessor(),
            ContentType.DOCUMENT: DocumentProcessor()
        }
    
    async def detect_content_type(self, file_path: str) -> ContentType:
        """Detect content type from file."""
        mime_type, _ = mimetypes.guess_type(file_path)
        
        if not mime_type:
            # Fallback to file extension
            ext = Path(file_path).suffix.lower()
            mime_type = {
                '.txt': 'text/plain',
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg', 
                '.png': 'image/png',
                '.gif': 'image/gif',
                '.mp3': 'audio/mpeg',
                '.wav': 'audio/wav',
                '.mp4': 'video/mp4',
                '.avi': 'video/avi',
                '.pdf': 'application/pdf',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            }.get(ext, 'application/octet-stream')
        
        # Map MIME types to content types
        if mime_type.startswith('text/'):
            return ContentType.TEXT
        elif mime_type.startswith('image/'):
            return ContentType.IMAGE
        elif mime_type.startswith('audio/'):
            return ContentType.AUDIO
        elif mime_type.startswith('video/'):
            return ContentType.VIDEO
        elif mime_type in ['application/pdf', 'application/msword'] or 'document' in mime_type:
            return ContentType.DOCUMENT
        else:
            return ContentType.TEXT  # Default fallback
    
    async def process_file(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process a file with the appropriate processor."""
        try:
            # Detect content type
            content_type = await self.detect_content_type(file_path)
            mime_type, _ = mimetypes.guess_type(file_path)
            
            # Get appropriate processor
            processor = self.processors[content_type]
            
            # Check if processor can handle this file
            if not await processor.can_process(file_path, mime_type or ''):
                raise ProcessingError(f"No processor available for {mime_type}")
            
            # Process the file
            result = await processor.process(file_path, **kwargs)
            
            # Add common metadata
            result.update({
                'content_type': content_type,
                'mime_type': mime_type,
                'processing_status': ProcessingStatus.COMPLETED
            })
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            raise ProcessingError(f"Processing failed: {e}")


# Export main service
__all__ = [
    'MultiModalProcessingService',
    'ProcessingError',
    'TextProcessor',
    'ImageProcessor', 
    'AudioProcessor',
    'VideoProcessor',
    'DocumentProcessor'
]
